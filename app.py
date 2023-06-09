import os
import openai
import json
import re
import fontawesome as fa
import tempfile
import boto3
from botocore.exceptions import NoCredentialsError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.util import Pt
import streamlit as st

fonts_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")
os.environ["FONTCONFIG_PATH"] = fonts_path

openai.api_key = os.environ["OPENAI_API_KEY"]

def apply_theme(presentation, theme):
	if theme == "dark":
		background_color = RGBColor(43, 43, 43)
		text_color = RGBColor(255, 255, 255)
		font_name = "Calibri"
	elif theme == "light":
		background_color = RGBColor(239, 239, 239)
		text_color = RGBColor(32, 32, 32)
		font_name = "Calibri"
	elif theme == "corporate":
		background_color = RGBColor(46, 117, 182)
		text_color = RGBColor(255, 255, 255)
		font_name = "Arial"
	elif theme == "playful":
		background_color = RGBColor(255, 204, 102)
		text_color = RGBColor(32, 32, 32)
		font_name = "Comic Sans MS"
	elif theme == "modern":
		background_color = RGBColor(45, 62, 80)
		text_color = RGBColor(255, 255, 255)
		font_name = "Segoe UI"
	elif theme == "vibrant":
		background_color = RGBColor(236, 98, 128)
		text_color = RGBColor(255, 255, 255)
		font_name = "Verdana"
	else:
		# Default theme
		background_color = RGBColor(239, 239, 239)
		text_color = RGBColor(32, 32, 32)
		font_name = "Calibri"

	for slide in presentation.slides:
		slide.background.fill.solid()
		slide.background.fill.fore_color.rgb = background_color
		for shape in slide.shapes:
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					for run in paragraph.runs:
						run.font.color.rgb = text_color
						run.font.name = font_name
						if theme in ["corporate", "modern", "vibrant"]:
							run.font.bold = True

	return presentation

def upload_to_s3_and_get_temporary_url(bucket_name, file_path, file_key, expiration=3600):
	aws_access_key_id = os.environ.get("AWS_ACCESS_KEY_ID")
	aws_secret_access_key = os.environ.get("AWS_SECRET_ACCESS_KEY")
	
	s3 = boto3.client(
		"s3",
		aws_access_key_id=aws_access_key_id,
		aws_secret_access_key=aws_secret_access_key,
	)
	
	try:
		s3.upload_file(file_path, bucket_name, file_key)
		response = s3.generate_presigned_url(
			"get_object",
			Params={"Bucket": bucket_name, "Key": file_key},
			ExpiresIn=expiration,
		)
		return response
	except NoCredentialsError as e:
		raise BadRequestError(f"Credentials not found: {str(e)}")

def generate_pptx(lesson_topic):
	prompt = (
		"Create PowerPoint slides for a lesson plan. The slides should be visually engaging, include concise headings and bullet points, and have relevant images or icons when necessary. Limit each slide to a maximum of 4 sub-points and a single image or icon when relevant. Divide the same heading into multiple slides if required to make the points more clear."
		"\n\nFor the first slide, include the lesson title and relevant sub-points. Also, include a closing slide with takeaways from the lesson. Choose a PowerPoint theme from these options: dark, light, corporate, playful, modern, and vibrant, depending on the lesson's context."
		"\n\nThe output should be suitable for use with the python-pptx library to create a PowerPoint presentation."
		"\n\nLesson Plan:\n{lesson_topic}"
		"\n\nFor each slide, provide this information:\n\n"
		"#. Slide (slide_title):\n"
		"Heading: concise_heading\n"
		"Sub-point 1:\n"
		"Sub-point 2:\n"
		"...\n"
		"If an image is relevant, include: 'Image: short_description_of_image'\n"
		"If an icon is relevant, include: 'Icon: font_awesome_icon_code'\n"
		"When creating the slides, remember to use clear and concise language, write the slides for the students to understand, and use appropriate images or icons, and choose a suitable theme for the PowerPoint presentation."
	)

	full_prompt = "".join(prompt).format(lesson_topic=lesson_topic)

	response = openai.ChatCompletion.create(
	model="gpt-4",
	messages=[
		{
			"role": "system",
			"content": (
					"You are a helpful assistant capable of creating clear and concise PowerPoint slide outlines used by teachers during their lessons based on a given lesson plan. You follow template instructions carefully"
			),
		},
		{"role": "user", "content": full_prompt},
	],
	max_tokens=650,
	n=1,
	stop=None,
	temperature=0.8,
	#    top_p=0.9,
	)
	output = response.choices[0].message['content'].strip()

	theme_pattern = re.compile(r"(dark|light|corporate|playful|modern|vibrant)")
	theme_match = theme_pattern.search(output)
	if theme_match:
		theme = theme_match.group(0)
	else:
		theme = "default"  # or any default theme you want to use

	output = theme_pattern.sub("", output).strip()

	ppt = Presentation()

	slides_data = output.split('\n\n')

	for slide_data in slides_data:
		slide_info = slide_data.split('\n')

		slide = ppt.slides.add_slide(ppt.slide_layouts[1])

		title = slide.shapes.title
		title.text = slide_info[0].split(':', 1)[-1].strip()

		content = slide_info[1:]
		content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4))
		content_text = content_box.text_frame

		placeholder_shape = None

		for line in content:
			if re.match(r"Image Placeholder|Icon:", line):
				image_placeholder_left = Inches(4.5)
				placeholder_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, image_placeholder_left, Inches(1.5), Inches(2), Inches(2))
				if "Icon" in line:
					icon_code = line.split("Icon:")[-1].strip()
					icon = fa.icons.get(icon_code)
					
					if icon is not None:
						placeholder_shape.text_frame.text = icon
						placeholder_shape.text_frame.paragraphs[0].runs[0].font.name = "Font Awesome 5 Free"
						placeholder_shape.text_frame.paragraphs[0].runs[0].font.bold = True
						placeholder_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(72)
					else:
						placeholder_shape.text_frame.text = f"Icon not found: {icon_code}"

		apply_theme(ppt, theme)

		with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt_file:
			ppt.save(tmp_ppt_file.name)
			temp_file_path = tmp_ppt_file.name

			# Extract the lesson title from the first slide
			first_slide_data = slides_data[0]
			first_slide_title_line = first_slide_data.split('\n')[0]
			
			if "Slide" in first_slide_title_line:
				first_slide_title = first_slide_title_line.split(':', 1)[-1].strip()
			else:
				first_slide_title = "PowerPoint_from_Lesson_Plan"
				
			s3_bucket_name = os.environ["S3_BUCKET_NAME"]
			file_key = f"{first_slide_title.replace(' ', '_')}_presentation.pptx"
			presigned_url = upload_to_s3_and_get_temporary_url(s3_bucket_name, temp_file_path, file_key, expiration=3600)

	return {'temporary_url': presigned_url}

st.title("PowerPoint Presentation Generator")
lesson_topic = st.text_area("Enter the lesson topic:")

if st.button("Generate Presentation"):
	if lesson_topic:
		result = generate_pptx(lesson_topic)
		st.write(f"Your presentation is ready! [Download it here]({result['temporary_url']})")
	else:
		st.warning("Please enter a lesson topic.")
